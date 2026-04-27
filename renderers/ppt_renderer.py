"""
ppt_renderer.py
---------------
Renders PowerPoint templates using python-pptx.
Templates: discovery_call, onboarding_summary
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
C_DARK   = RGBColor(0x0D, 0x11, 0x17)
C_ACCENT = RGBColor(0x25, 0x63, 0xEB)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)
C_MUTED  = RGBColor(0x64, 0x74, 0x8B)
C_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────
def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


def _fill_shape(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _add_rect(slide, left, top, width, height, rgb: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _fill_shape(shape, rgb)
    shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              font_name="Calibri", font_size=18, bold=False,
              color=C_DARK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _bullet_frame(slide, items: list, left, top, width, height,
                  font_size=16, color=C_DARK, marker="●"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{marker}  {item}"
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
#  TEMPLATE 1 — DISCOVERY CALL BRIEF
# ══════════════════════════════════════════════════════════════════════════════
def render_discovery_call(data: dict, output_path: str):
    client    = data.get("client_name", "Prospective Client")
    date      = data.get("meeting_date", "Upcoming")
    prep_by   = data.get("prepared_by", "Agentic Solutions")
    agenda    = data.get("agenda", [])
    pain_pts  = data.get("pain_points", [])
    strengths = data.get("our_strengths", [])

    prs = _prs()

    # ── SLIDE 1: Title ──────────────────────────────────────────────────────
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_DARK)
    _add_rect(s, 0, Inches(5.8), SLIDE_W, Inches(0.15), C_ACCENT)

    _add_text(s, "DISCOVERY CALL BRIEF",
              Inches(1), Inches(1.6), Inches(11), Inches(1),
              font_size=38, bold=True, color=C_WHITE)

    _add_text(s, client,
              Inches(1), Inches(2.8), Inches(9), Inches(0.8),
              font_size=22, color=RGBColor(0xCB, 0xD5, 0xE1))

    _add_text(s, f"Meeting Date: {date}   ·   Prepared by: {prep_by}",
              Inches(1), Inches(3.5), Inches(10), Inches(0.5),
              font_size=13, color=RGBColor(0x94, 0xA3, 0xB8))

    # Decorative accent bar
    _add_rect(s, Inches(1), Inches(2.5), Inches(1.5), Inches(0.06), C_ACCENT)

    # ── SLIDE 2: Agenda ─────────────────────────────────────────────────────
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_LIGHT)
    _add_rect(s, 0, 0, Inches(0.08), SLIDE_H, C_ACCENT)
    _add_rect(s, 0, 0, SLIDE_W, Inches(1.4), C_DARK)

    _add_text(s, "Today's Agenda",
              Inches(0.4), Inches(0.25), Inches(10), Inches(0.9),
              font_size=26, bold=True, color=C_WHITE)

    for i, item in enumerate(agenda):
        top = Inches(1.7 + i * 0.9)
        _add_rect(s, Inches(0.4), top, Inches(0.5), Inches(0.5), C_ACCENT)
        _add_text(s, str(i+1), Inches(0.4), top, Inches(0.5), Inches(0.5),
                  font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, item, Inches(1.1), top + Inches(0.05), Inches(11), Inches(0.5),
                  font_size=16, color=C_DARK)

    # ── SLIDE 3: Pain Points ────────────────────────────────────────────────
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_LIGHT)
    _add_rect(s, 0, 0, Inches(0.08), SLIDE_H, C_ACCENT)
    _add_rect(s, 0, 0, SLIDE_W, Inches(1.4), C_DARK)

    _add_text(s, "Understanding Your Challenges",
              Inches(0.4), Inches(0.25), Inches(12), Inches(0.9),
              font_size=26, bold=True, color=C_WHITE)

    _add_text(s, "We've heard these challenges from organisations like yours:",
              Inches(0.4), Inches(1.6), Inches(12), Inches(0.5),
              font_size=13, italic=True, color=C_MUTED)

    for i, pt in enumerate(pain_pts):
        top = Inches(2.3 + i * 1.1)
        _add_rect(s, Inches(0.4), top, Inches(11.5), Inches(0.75), C_WHITE)
        _add_text(s, f"⚡  {pt}", Inches(0.6), top + Inches(0.12), Inches(11), Inches(0.5),
                  font_size=15, color=C_DARK)

    # ── SLIDE 4: Our Strengths ──────────────────────────────────────────────
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_DARK)
    _add_rect(s, 0, 0, Inches(0.08), SLIDE_H, C_ACCENT)

    _add_text(s, "Why Partner With Us",
              Inches(0.4), Inches(0.4), Inches(12), Inches(0.9),
              font_size=26, bold=True, color=C_WHITE)
    _add_rect(s, Inches(0.4), Inches(1.3), Inches(2), Inches(0.06), C_ACCENT)

    cols = min(len(strengths), 3)
    box_w = Inches(3.8)
    for i, st in enumerate(strengths[:cols]):
        left = Inches(0.4 + i * 4.2)
        _add_rect(s, left, Inches(1.7), box_w, Inches(3.5), RGBColor(0x1E, 0x29, 0x3B))
        _add_rect(s, left, Inches(1.7), box_w, Inches(0.1), C_ACCENT)
        _add_text(s, "✓", left + Inches(0.2), Inches(2.0), Inches(0.5), Inches(0.5),
                  font_size=22, bold=True, color=C_ACCENT)
        _add_text(s, st, left + Inches(0.2), Inches(2.6), box_w - Inches(0.4), Inches(2),
                  font_size=14, color=C_WHITE)

    # ── SLIDE 5: Next Steps ─────────────────────────────────────────────────
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_LIGHT)
    _add_rect(s, 0, 0, SLIDE_W, Inches(1.4), C_ACCENT)
    _add_text(s, "Next Steps",
              Inches(0.4), Inches(0.25), Inches(10), Inches(0.9),
              font_size=26, bold=True, color=C_WHITE)

    steps = [
        "Share current tooling & integration overview",
        "Technical discovery call with your IT team",
        "Receive tailored proposal within 3 business days",
        "Kick off pilot engagement",
    ]
    for i, step in enumerate(steps):
        top = Inches(1.7 + i * 1.1)
        _add_rect(s, Inches(0.4), top, Inches(0.6), Inches(0.6), C_DARK)
        _add_text(s, str(i+1), Inches(0.4), top, Inches(0.6), Inches(0.6),
                  font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, step, Inches(1.2), top + Inches(0.1), Inches(11), Inches(0.5),
                  font_size=16, color=C_DARK)

    prs.save(output_path)
    return output_path


# ══════════════════════════════════════════════════════════════════════════════
#  TEMPLATE 2 — ONBOARDING SUMMARY PPT (stretch / second PPT template)
# ══════════════════════════════════════════════════════════════════════════════
def render_onboarding_summary(data: dict, output_path: str):
    client = data.get("client_name", "Your Company")
    fields = data.get("rep_fields", ["Full Name", "Role", "Start Date", "Manager"])

    prs = _prs()

    # Title slide
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_ACCENT)
    _add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.1), C_WHITE)
    _add_text(s, "Sales Rep Onboarding", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
              font_size=36, bold=True, color=C_WHITE)
    _add_text(s, client, Inches(1), Inches(3), Inches(9), Inches(0.7),
              font_size=20, color=RGBColor(0xBF, 0xDB, 0xFE))
    _add_text(s, "Your guide to a great first week.",
              Inches(1), Inches(3.8), Inches(10), Inches(0.6),
              font_size=14, italic=True, color=C_WHITE)

    # Fields overview slide
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_LIGHT)
    _add_rect(s, 0, 0, SLIDE_W, Inches(1.4), C_DARK)
    _add_text(s, "Required Onboarding Information",
              Inches(0.4), Inches(0.25), Inches(12), Inches(0.9),
              font_size=24, bold=True, color=C_WHITE)

    for i, field in enumerate(fields):
        top = Inches(1.7 + i * 0.95)
        _add_rect(s, Inches(0.4), top, Inches(12), Inches(0.7), C_WHITE)
        _add_rect(s, Inches(0.4), top, Inches(0.08), Inches(0.7), C_ACCENT)
        _add_text(s, field, Inches(0.65), top + Inches(0.12), Inches(10), Inches(0.5),
                  font_size=15, color=C_DARK)

    prs.save(output_path)
    return output_path


# ── Dispatcher ──────────────────────────────────────────────────────────────
def render_ppt(template_key: str, data: dict, output_path: str) -> str:
    if template_key == "discovery_call":
        return render_discovery_call(data, output_path)
    else:
        return render_onboarding_summary(data, output_path)
